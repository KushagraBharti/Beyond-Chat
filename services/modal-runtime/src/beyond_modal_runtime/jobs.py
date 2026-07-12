"""Deterministic document and data/finance output jobs executed inside a sandbox."""

from __future__ import annotations

import argparse
import csv
import hashlib
import json
import subprocess
from pathlib import Path
from typing import Any


def digest(path: Path) -> str:
    return f"sha256:{hashlib.sha256(path.read_bytes()).hexdigest()}"


def create_presentation(path: Path, title: str, bullets: list[str]) -> dict[str, Any]:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    path.parent.mkdir(parents=True, exist_ok=True)
    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(1.0))
    title_paragraph = title_box.text_frame.paragraphs[0]
    title_paragraph.text = title
    title_paragraph.font.size = Pt(34)
    title_paragraph.font.bold = True
    body = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.0), Inches(4.5)).text_frame
    body.clear()
    for index, bullet in enumerate(bullets):
        paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(22)
    deck.core_properties.title = title
    deck.core_properties.author = "Beyond Chat"
    deck.save(path)
    return {"path": str(path), "digest": digest(path), "slides": len(deck.slides)}


def edit_presentation(source: Path, target: Path, suffix: str) -> dict[str, Any]:
    from pptx import Presentation

    deck = Presentation(source)
    for slide in deck.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                shape.text_frame.paragraphs[0].text = f"{shape.text_frame.paragraphs[0].text}{suffix}"
                deck.save(target)
                return {"path": str(target), "digest": digest(target), "edited": True}
    raise ValueError("presentation contains no editable text")


def render_presentation(source: Path, output_directory: Path) -> dict[str, Any]:
    output_directory.mkdir(parents=True, exist_ok=True)
    process = subprocess.run(
        ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", str(output_directory), str(source)],
        check=False,
        capture_output=True,
        text=True,
        timeout=120,
    )
    if process.returncode != 0:
        raise RuntimeError(f"LibreOffice render failed with exit code {process.returncode}")
    target = output_directory / f"{source.stem}.pdf"
    if not target.exists() or target.stat().st_size == 0:
        raise RuntimeError("LibreOffice did not produce a non-empty PDF")
    return {"path": str(target), "digest": digest(target), "byte_size": target.stat().st_size}


def finance_summary(source: Path, output: Path) -> dict[str, Any]:
    with source.open("r", encoding="utf-8", newline="") as stream:
        rows = list(csv.DictReader(stream))
    amounts = [float(row["amount"]) for row in rows]
    payload = {
        "count": len(amounts),
        "sum": round(sum(amounts), 4),
        "average": round(sum(amounts) / len(amounts), 4) if amounts else 0,
        "minimum": min(amounts) if amounts else None,
        "maximum": max(amounts) if amounts else None,
    }
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text(json.dumps(payload, sort_keys=True, separators=(",", ":")), encoding="utf-8")
    return {"path": str(output), "digest": digest(output), **payload}


def main() -> None:
    parser = argparse.ArgumentParser()
    subcommands = parser.add_subparsers(dest="command", required=True)
    presentation = subcommands.add_parser("presentation")
    presentation.add_argument("output", type=Path)
    presentation.add_argument("--title", required=True)
    presentation.add_argument("--bullet", action="append", default=[])
    edit = subcommands.add_parser("edit-presentation")
    edit.add_argument("source", type=Path)
    edit.add_argument("output", type=Path)
    edit.add_argument("--suffix", default=" — revised")
    render = subcommands.add_parser("render-presentation")
    render.add_argument("source", type=Path)
    render.add_argument("output_directory", type=Path)
    finance = subcommands.add_parser("finance-summary")
    finance.add_argument("source", type=Path)
    finance.add_argument("output", type=Path)
    args = parser.parse_args()
    if args.command == "presentation":
        result = create_presentation(args.output, args.title, args.bullet)
    elif args.command == "edit-presentation":
        result = edit_presentation(args.source, args.output, args.suffix)
    elif args.command == "render-presentation":
        result = render_presentation(args.source, args.output_directory)
    else:
        result = finance_summary(args.source, args.output)
    print(json.dumps(result, sort_keys=True, separators=(",", ":")))


if __name__ == "__main__":
    main()
